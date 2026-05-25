#!/usr/bin/env python3
"""14 页项目展示 PPT：三部分结构 + 总目录 + 分章目录 + 首尾。"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
ASSETS = ROOT / "output" / "presentation_assets"
OUT = ROOT / "output" / "文档智能系统_项目展示_14页.pptx"

INK = RGBColor(0x11, 0x18, 0x27)
MUTE = RGBColor(0x64, 0x74, 0x8B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
NAVY = RGBColor(0x0B, 0x12, 0x20)
TEAL = RGBColor(0x14, 0xB8, 0xA6)
BLUE = RGBColor(0x25, 0x63, 0xEB)
BG = RGBColor(0xF1, 0xF5, 0xF9)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
FONT = "Microsoft YaHei"
TOTAL = 14


def asset(name: str) -> Path:
    p = ASSETS / name
    if not p.is_file():
        raise FileNotFoundError(f"缺少截图: {p}")
    return p


def blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def fill_bg(slide, color: RGBColor) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text,
    *,
    size=18,
    bold=False,
    color=INK,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = FONT
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return box


def add_bullets(slide, left, top, width, height, lines: list[str], *, size=20, color=INK):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.name = FONT
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(10)
        if line.startswith("•") or line.startswith("-"):
            p.bullet = True


def add_picture_fit(slide, path: Path, left, top, max_w, max_h, caption: str | None = None):
    pic = slide.shapes.add_picture(str(path), left, top)
    ratio = min(max_w / pic.width, max_h / pic.height, 1.0)
    pic.width = int(pic.width * ratio)
    pic.height = int(pic.height * ratio)
    if caption:
        add_textbox(slide, left, top + pic.height + Inches(0.06), max_w, Inches(0.32), caption, size=13, color=MUTE)
    return pic


def footer(slide, n: int):
    add_textbox(slide, Inches(0.55), Inches(7.05), Inches(5), Inches(0.28), f"文档智能系统  ·  {n}/{TOTAL}", size=11, color=MUTE)


def header_content(slide, part: str, title: str):
    add_textbox(slide, Inches(0.55), Inches(0.35), Inches(4), Inches(0.28), part, size=13, bold=True, color=TEAL)
    add_textbox(slide, Inches(0.55), Inches(0.68), Inches(12), Inches(0.62), title, size=30, bold=True, color=INK)
    bar = slide.shapes.add_shape(1, Inches(0.55), Inches(1.28), Inches(0.95), Inches(0.05))
    bar.fill.solid()
    bar.fill.fore_color.rgb = TEAL
    bar.line.fill.background()


def slide_cover(prs: Presentation):
    s = blank_slide(prs)
    fill_bg(s, NAVY)
    add_rect_accent(s)
    add_textbox(s, Inches(0.75), Inches(0.6), Inches(6), Inches(0.35), "DOCUMENT INTELLIGENCE", size=14, bold=True, color=TEAL)
    add_textbox(s, Inches(0.75), Inches(1.15), Inches(6.5), Inches(1.0), "文档智能系统", size=50, bold=True, color=WHITE)
    add_textbox(
        s,
        Inches(0.78),
        Inches(2.35),
        Inches(5.8),
        Inches(1.2),
        "基于大语言模型的文档理解与多源数据融合\nElectron 桌面版 · 项目展示",
        size=20,
        color=RGBColor(0xCB, 0xD5, 0xE1),
    )
    add_picture_fit(s, asset("img06.png"), Inches(6.4), Inches(0.5), Inches(6.4), Inches(6.3))


def add_rect_accent(s, dark=True):
    if dark:
        s.shapes.add_shape(1, Inches(11.2), Inches(-0.3), Inches(2.8), Inches(8.2)).fill.solid()
        r = s.shapes[-1]
        r.fill.fore_color.rgb = RGBColor(0x0F, 0x76, 0x6E)
        r.line.fill.background()
        s.shapes.add_shape(1, Inches(12.0), Inches(-0.5), Inches(2.0), Inches(8.5)).fill.solid()
        r2 = s.shapes[-1]
        r2.fill.fore_color.rgb = RGBColor(0x1D, 0x4E, 0xD8)
        r2.line.fill.background()


def slide_master_toc(prs: Presentation):
    """总目录"""
    s = blank_slide(prs)
    fill_bg(s, BG)
    add_textbox(s, Inches(0.75), Inches(0.55), Inches(4), Inches(0.4), "目录", size=16, bold=True, color=TEAL)
    add_textbox(s, Inches(0.75), Inches(1.0), Inches(8), Inches(0.8), "汇报提纲", size=40, bold=True, color=INK)
    parts = [
        ("01", "项目定位与架构", "定位 · 产品入口 · 架构与智能体", "第 3～6 页"),
        ("02", "核心功能展示", "文档库 · 智能对话 · 工作流 · 实测结果", "第 7～11 页"),
        ("03", "进度与规划", "阶段成果 · 规划展望", "第 12～14 页"),
    ]
    y = Inches(2.0)
    for num, title, sub, pages in parts:
        card = s.shapes.add_shape(1, Inches(0.75), y, Inches(11.8), Inches(1.35))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
        add_textbox(s, Inches(1.0), y + Inches(0.22), Inches(0.7), Inches(0.55), num, size=28, bold=True, color=BLUE)
        add_textbox(s, Inches(1.85), y + Inches(0.18), Inches(7.5), Inches(0.5), title, size=24, bold=True, color=INK)
        add_textbox(s, Inches(1.85), y + Inches(0.68), Inches(7.5), Inches(0.4), sub, size=15, color=MUTE)
        add_textbox(s, Inches(10.2), y + Inches(0.4), Inches(1.8), Inches(0.4), pages, size=14, color=TEAL, align=PP_ALIGN.RIGHT)
        y += Inches(1.55)
    footer(s, 2)


def slide_part_toc(prs: Presentation, n: int, part_no: str, part_title: str, items: list[tuple[str, str]]):
    """分章目录（每部分一页）"""
    s = blank_slide(prs)
    fill_bg(s, NAVY if part_no == "01" else (RGBColor(0x0F, 0x17, 0x2A) if part_no == "02" else RGBColor(0x1E, 0x29, 0x3B)))
    accent = TEAL if part_no == "01" else (BLUE if part_no == "02" else AMBER)
    add_textbox(s, Inches(0.75), Inches(0.55), Inches(2), Inches(0.35), f"PART {part_no}", size=14, bold=True, color=accent)
    add_textbox(s, Inches(0.75), Inches(1.05), Inches(11), Inches(0.9), part_title, size=38, bold=True, color=WHITE)
    add_textbox(s, Inches(0.75), Inches(2.0), Inches(3), Inches(0.35), "本章目录", size=16, color=RGBColor(0x94, 0xA3, 0xB8))
    y = Inches(2.55)
    for i, (title, desc) in enumerate(items, 1):
        add_textbox(s, Inches(0.95), y, Inches(0.5), Inches(0.4), f"{i}.", size=20, bold=True, color=accent)
        add_textbox(s, Inches(1.45), y, Inches(9.5), Inches(0.42), title, size=22, bold=True, color=WHITE)
        add_textbox(s, Inches(1.45), y + Inches(0.42), Inches(9.5), Inches(0.38), desc, size=16, color=RGBColor(0xCB, 0xD5, 0xE1))
        y += Inches(1.05)
    footer(s, n)


def slide_positioning(prs: Presentation, n: int):
    s = blank_slide(prs)
    fill_bg(s, BG)
    header_content(s, "第一部分", "项目定位：解决什么问题")
    add_bullets(
        s,
        Inches(0.6),
        Inches(1.55),
        Inches(5.4),
        Inches(4.8),
        [
            "企业积累大量非结构化文档（公报、合同、纪要、报表）",
            "格式不一、来源分散，人工阅读与录入成本高",
            "理解、翻译、摘录、填表分散在多个工具，难形成稳定流程",
            "",
            "定位：基于 LLM 的本地桌面工作台",
            "读文档 → 理解/抽取 → 写入模板 → 另存为",
        ],
        size=19,
    )
    add_picture_fit(s, asset("img01.png"), Inches(6.1), Inches(1.45), Inches(6.6), Inches(5.2), "文档理解：多文件讲解与问答")
    footer(s, n)


def slide_product_entries(prs: Presentation, n: int):
    s = blank_slide(prs)
    fill_bg(s, BG)
    header_content(s, "第一部分", "产品入口：资料 · 对话 · 工作流")
    entries = [
        ("文档库", "上传与管理 PDF / Word / Excel / Markdown"),
        ("智能对话", "默认对话 · 文档理解 · 文档编辑 · 提取与填表"),
        ("工作流", "输入 → AI（翻译/摘要/抽取）→ 输出文件"),
    ]
    y = Inches(1.65)
    for title, desc in entries:
        card = s.shapes.add_shape(1, Inches(0.6), y, Inches(5.8), Inches(1.15))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
        add_textbox(s, Inches(0.85), y + Inches(0.15), Inches(5.2), Inches(0.4), title, size=22, bold=True, color=BLUE)
        add_textbox(s, Inches(0.85), y + Inches(0.55), Inches(5.2), Inches(0.45), desc, size=16, color=MUTE)
        y += Inches(1.35)
    add_textbox(
        s,
        Inches(0.6),
        Inches(5.55),
        Inches(5.8),
        Inches(0.7),
        "生成结果通过系统「另存为」保存，不强制写入固定目录。",
        size=15,
        bold=True,
        color=BLUE,
    )
    add_picture_fit(s, asset("img11.png"), Inches(6.5), Inches(1.5), Inches(6.2), Inches(5.2), "主导航：文档库 / 智能对话 / 工作流")
    footer(s, n)


def slide_architecture(prs: Presentation, n: int):
    s = blank_slide(prs)
    fill_bg(s, BG)
    header_content(s, "第一部分", "技术架构与三类智能体（本地）")
    layers = [
        ("桌面壳", "Electron：窗口、另存为、托管本地服务"),
        ("交互层", "文档库 · 智能对话 · 工作流画布 · 模型设置"),
        ("本地服务层", "会话与文件 · 流式对话 · 工作流执行"),
        ("智能处理层", "文档解析 · 三类智能体 · 大模型调用"),
        ("本地数据层", "配置、文档库、会话、工作流（无需数据库）"),
    ]
    y = Inches(1.48)
    for name, desc in layers:
        bar = s.shapes.add_shape(1, Inches(0.55), y, Inches(0.12), Inches(0.72))
        bar.fill.solid()
        bar.fill.fore_color.rgb = TEAL
        bar.line.fill.background()
        add_textbox(s, Inches(0.78), y + Inches(0.06), Inches(1.45), Inches(0.32), name, size=15, bold=True, color=BLUE)
        add_textbox(s, Inches(2.3), y + Inches(0.06), Inches(4.2), Inches(0.5), desc, size=14, color=INK)
        y += Inches(0.82)
    add_textbox(s, Inches(6.35), Inches(1.48), Inches(2.2), Inches(0.38), "智能体 A", size=17, bold=True, color=TEAL)
    add_textbox(s, Inches(6.35), Inches(1.82), Inches(3.2), Inches(0.55), "文档理解、自然语言编辑", size=14, color=INK)
    add_textbox(s, Inches(6.35), Inches(2.45), Inches(2.2), Inches(0.38), "智能体 B", size=17, bold=True, color=BLUE)
    add_textbox(s, Inches(6.35), Inches(2.78), Inches(3.2), Inches(0.55), "实体/字段抽取 → 结构化数据", size=14, color=INK)
    add_textbox(s, Inches(6.35), Inches(3.42), Inches(2.2), Inches(0.38), "智能体 D", size=17, bold=True, color=AMBER)
    add_textbox(s, Inches(6.35), Inches(3.75), Inches(3.2), Inches(0.55), "Excel 筛选并填入 Word/Excel 模板", size=14, color=INK)
    add_picture_fit(s, asset("img02.png"), Inches(9.5), Inches(1.45), Inches(3.35), Inches(5.25), "结构化抽取与预览")
    add_picture_fit(s, asset("img14.png"), Inches(6.35), Inches(4.55), Inches(2.9), Inches(1.95))
    footer(s, n)


def slide_agents(prs: Presentation, n: int):
    s = blank_slide(prs)
    fill_bg(s, BG)
    header_content(s, "第一部分", "三类智能体分工（本地）")
    rows = [
        ("智能体 A", "文档理解、自然语言编辑", TEAL),
        ("智能体 B", "实体/字段抽取 → 结构化数据", BLUE),
        ("智能体 D", "Excel 筛选并填入 Word/Excel 模板", AMBER),
    ]
    y = Inches(1.55)
    for name, desc, accent in rows:
        add_textbox(s, Inches(0.6), y, Inches(2.2), Inches(0.42), name, size=22, bold=True, color=accent)
        add_textbox(s, Inches(0.6), y + Inches(0.45), Inches(5.5), Inches(0.45), desc, size=18, color=INK)
        y += Inches(1.15)
    add_picture_fit(s, asset("img02.png"), Inches(6.3), Inches(1.45), Inches(6.5), Inches(5.0), "结构化抽取与预览")
    footer(s, n)


def slide_library(prs: Presentation, n: int):
    s = blank_slide(prs)
    fill_bg(s, BG)
    header_content(s, "第二部分", "文档库：多格式资料统一管理")
    add_picture_fit(s, asset("img04.png"), Inches(0.55), Inches(1.45), Inches(12.2), Inches(5.35), "按空间管理，供对话与工作流选用")
    footer(s, n)


def slide_chat(prs: Presentation, n: int):
    s = blank_slide(prs)
    fill_bg(s, BG)
    header_content(s, "第二部分", "智能对话：理解 · 编辑 · 填表")
    add_picture_fit(s, asset("img09.png"), Inches(0.55), Inches(1.45), Inches(7.6), Inches(5.35), "基于已选文件流式回答")
    add_picture_fit(s, asset("img03.png"), Inches(8.4), Inches(1.45), Inches(4.3), Inches(2.45), "生成结果另存为")
    add_picture_fit(s, asset("img07.png"), Inches(8.4), Inches(4.05), Inches(4.3), Inches(2.75), "多厂商模型配置")
    footer(s, n)


def slide_workflow(prs: Presentation, n: int):
    s = blank_slide(prs)
    fill_bg(s, BG)
    header_content(s, "第二部分", "工作流：可复用的批处理")
    add_picture_fit(s, asset("img06.png"), Inches(0.55), Inches(1.45), Inches(8.8), Inches(5.35), "文档输入 → AI 翻译/摘要 → 文件输出")
    add_textbox(
        s,
        Inches(9.5),
        Inches(1.7),
        Inches(3.3),
        Inches(4.5),
        "内置示例\n· 文档翻译流\n· 内容提取流\n\n节点可配置目标语言、\n输出 PDF / Markdown",
        size=17,
        color=INK,
    )
    footer(s, n)


def slide_validation(prs: Presentation, n: int):
    s = blank_slide(prs)
    fill_bg(s, BG)
    header_content(s, "第二部分", "实测：填表 · 编辑 · 报告导出")
    imgs = [
        ("img08.png", "Excel 数据源"),
        ("img13.png", "填表结果"),
        ("img12.png", "Word 输出"),
        ("img10.png", "PDF 报告"),
    ]
    x0, y0 = Inches(0.55), Inches(1.48)
    w, h = Inches(5.9), Inches(2.42)
    gap = Inches(0.22)
    for i, (name, cap) in enumerate(imgs):
        col, row = i % 2, i // 2
        add_picture_fit(s, asset(name), x0 + col * (w + gap), y0 + row * (h + Inches(0.42)), w, h - Inches(0.32), cap)
    footer(s, n)


def slide_progress(prs: Presentation, n: int):
    s = blank_slide(prs)
    fill_bg(s, BG)
    header_content(s, "第三部分", "阶段成果")
    done = [
        "Electron 桌面三栏：文档库 / 智能对话 / 工作流",
        "三类智能体 + 工作流编排，覆盖理解、编辑、抽取、填表",
        "本地数据与模型设置，Windows 安装包可交付",
        "21+ 测试样例文档，覆盖主要办公格式",
        "生成文件支持系统「另存为」",
    ]
    add_bullets(s, Inches(0.6), Inches(1.55), Inches(5.6), Inches(5.2), [f"• {x}" for x in done], size=19)
    add_picture_fit(s, asset("img14.png"), Inches(6.4), Inches(1.5), Inches(6.4), Inches(4.8), "打包与本地运行")
    footer(s, n)


def slide_planning(prs: Presentation, n: int):
    s = blank_slide(prs)
    fill_bg(s, BG)
    header_content(s, "第三部分", "规划与展望")
    plans = [
        ("近期", "批量任务进度展示、失败重试与日志可读性"),
        ("中期", "更多行业模板（公报、采购、人事等）"),
        ("中期", "工作流节点扩展（更多 AI 能力组合）"),
        ("可选", "接入本地大模型，降低对外网 API 依赖"),
    ]
    y = Inches(1.55)
    for phase, content in plans:
        add_textbox(s, Inches(0.6), y, Inches(1.4), Inches(0.38), phase, size=20, bold=True, color=TEAL)
        add_textbox(s, Inches(2.1), y, Inches(10), Inches(0.55), content, size=19, color=INK)
        y += Inches(1.05)
    add_textbox(
        s,
        Inches(0.6),
        Inches(5.6),
        Inches(12),
        Inches(0.8),
        "目标：把「读文档、用文档」沉淀为可安装、可演示、可扩展的本地产品。",
        size=18,
        bold=True,
        color=BLUE,
    )
    footer(s, n)


def slide_end(prs: Presentation, n: int):
    s = blank_slide(prs)
    fill_bg(s, NAVY)
    add_rect_accent(s)
    add_textbox(s, Inches(0.85), Inches(2.4), Inches(10), Inches(1.0), "谢谢", size=54, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    add_textbox(
        s,
        Inches(0.88),
        Inches(3.55),
        Inches(10),
        Inches(1.2),
        "文档智能系统 · Electron 桌面版\n欢迎交流与提问",
        size=24,
        color=TEAL,
    )
    footer(s, n)


def slide_progress_and_plan(prs: Presentation, n: int):
    s = blank_slide(prs)
    fill_bg(s, BG)
    header_content(s, "第三部分", "阶段成果与规划")
    add_textbox(s, Inches(0.55), Inches(1.45), Inches(2.2), Inches(0.45), "已完成", size=22, bold=True, color=TEAL)
    add_bullets(
        s,
        Inches(0.55),
        Inches(1.95),
        Inches(5.6),
        Inches(3.2),
        [
            "• Electron 三栏：文档库 / 对话 / 工作流",
            "• 三类智能体 + 工作流，覆盖理解、编辑、抽取、填表",
            "• Windows 安装包与 21+ 测试样例",
            "• 生成文件系统「另存为」，本地数据无数据库",
        ],
        size=17,
    )
    add_textbox(s, Inches(6.35), Inches(1.45), Inches(2.2), Inches(0.45), "规划", size=22, bold=True, color=BLUE)
    add_bullets(
        s,
        Inches(6.35),
        Inches(1.95),
        Inches(5.8),
        Inches(2.8),
        [
            "• 近期：批量进度、失败重试与日志可读性",
            "• 中期：更多行业模板（公报、采购、人事）",
            "• 中期：工作流节点与 AI 能力组合扩展",
            "• 可选：接入本地大模型，降低外网依赖",
        ],
        size=17,
    )
    add_textbox(
        s,
        Inches(0.55),
        Inches(5.35),
        Inches(11.8),
        Inches(0.65),
        "目标：把「读文档、用文档」沉淀为可安装、可演示、可扩展的本地产品。",
        size=17,
        bold=True,
        color=BLUE,
    )
    add_picture_fit(s, asset("img14.png"), Inches(9.2), Inches(1.5), Inches(3.6), Inches(3.5), "打包与本地运行")
    footer(s, n)


def build() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_cover(prs)  # 1
    slide_master_toc(prs)  # 2

    slide_part_toc(
        prs,
        3,
        "01",
        "项目定位与架构",
        [
            ("项目定位与痛点", "非结构化文档处理的效率瓶颈"),
            ("产品三大入口", "文档库 · 智能对话 · 工作流"),
            ("架构与三类智能体", "本地五层 + 理解 / 抽取 / 填表"),
        ],
    )
    slide_positioning(prs, 4)
    slide_product_entries(prs, 5)
    slide_architecture(prs, 6)

    slide_part_toc(
        prs,
        7,
        "02",
        "核心功能展示",
        [
            ("文档库", "多格式资料上传与空间管理"),
            ("智能对话", "流式理解、编辑、另存为"),
            ("工作流", "翻译 / 摘要等批处理"),
            ("实测结果", "填表、Word、PDF 等输出验证"),
        ],
    )
    slide_library(prs, 8)
    slide_chat(prs, 9)
    slide_workflow(prs, 10)
    slide_validation(prs, 11)

    slide_part_toc(
        prs,
        12,
        "03",
        "进度与规划",
        [
            ("阶段成果", "已完成能力与交付物"),
            ("规划与展望", "近期与中期目标"),
        ],
    )
    slide_progress_and_plan(prs, 13)
    slide_end(prs, 14)

    if len(prs.slides) != TOTAL:
        raise RuntimeError(f"页数应为 {TOTAL}，实际 {len(prs.slides)}")

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    return OUT


def main() -> None:
    path = build()
    size_kb = path.stat().st_size // 1024
    print(f"已生成: {path} ({size_kb} KB, {TOTAL} 页)")


if __name__ == "__main__":
    main()
